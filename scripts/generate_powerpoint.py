#!/usr/bin/env python3
"""
Generate PowerPoint Report Template
Crée automatiquement un rapport PowerPoint avec 26 sections
"""

import os
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("❌ Module python-pptx non installé!")
    print("\n📦 Installation:")
    print("   pip3 install python-pptx")
    print("\nPuis réexécuter ce script.")
    exit(1)

# Configuration
OUTPUT_FILE = "../reports/Marketing_Performance_Report.pptx"
COMPANY_NAME = "Hulken"
REPORT_DATE = datetime.now().strftime("%B %Y")

# Colors (Brand)
COLOR_PRIMARY = RGBColor(66, 133, 244)  # Google Blue
COLOR_SUCCESS = RGBColor(52, 168, 83)   # Green
COLOR_WARNING = RGBColor(251, 188, 4)   # Yellow
COLOR_DANGER = RGBColor(234, 67, 53)    # Red
COLOR_DARK = RGBColor(60, 64, 67)       # Dark Gray
COLOR_LIGHT = RGBColor(241, 243, 244)   # Light Gray

def create_title_slide(prs, title, subtitle):
    """Crée une slide de titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title layout

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    return slide

def create_section_slide(prs, section_title):
    """Crée une slide de section"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section header
    title = slide.shapes.title
    title.text = section_title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    return slide

def create_content_slide(prs, title, bullet_points=None):
    """Crée une slide de contenu avec bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content

    title_shape = slide.shapes.title
    title_shape.text = title

    if bullet_points:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)

    return slide

def create_kpi_slide(prs, title, kpis):
    """Crée une slide avec des KPIs"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # KPIs (4 par ligne)
    kpi_width = Inches(2.2)
    kpi_height = Inches(1.5)
    start_x = Inches(0.5)
    start_y = Inches(1.8)
    spacing = Inches(0.2)

    for i, kpi in enumerate(kpis):
        row = i // 4
        col = i % 4

        x = start_x + col * (kpi_width + spacing)
        y = start_y + row * (kpi_height + spacing)

        # KPI Box
        kpi_box = slide.shapes.add_shape(
            1,  # Rectangle
            x, y, kpi_width, kpi_height
        )
        kpi_box.fill.solid()
        kpi_box.fill.fore_color.rgb = COLOR_LIGHT
        kpi_box.line.color.rgb = COLOR_PRIMARY

        # KPI Text
        text_frame = kpi_box.text_frame
        text_frame.clear()
        text_frame.margin_left = Pt(10)
        text_frame.margin_top = Pt(10)

        # Metric name
        p1 = text_frame.paragraphs[0]
        p1.text = kpi['name']
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_DARK

        # Value
        p2 = text_frame.add_paragraph()
        p2.text = kpi['value']
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_PRIMARY

        # YoY Change (if exists)
        if 'yoy' in kpi:
            p3 = text_frame.add_paragraph()
            yoy = float(kpi['yoy'].replace('%', ''))
            color = COLOR_SUCCESS if yoy >= 0 else COLOR_DANGER
            p3.text = f"↑ {kpi['yoy']}" if yoy >= 0 else f"↓ {kpi['yoy']}"
            p3.font.size = Pt(14)
            p3.font.color.rgb = color

    return slide

def create_table_slide(prs, title, headers, rows):
    """Crée une slide avec un tableau"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True

    # Table
    rows_count = len(rows) + 1  # +1 for header
    cols_count = len(headers)

    table = slide.shapes.add_table(
        rows_count, cols_count,
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(4.5)
    ).table

    # Header row
    for col_idx, header in enumerate(headers):
        cell = table.rows[0].cells[col_idx]
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.rows[row_idx + 1].cells[col_idx]
            cell.text = str(cell_data)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)

    return slide

def main():
    """Génère le PowerPoint complet"""

    print(f"\n🎨 Génération du rapport PowerPoint: {COMPANY_NAME}")
    print(f"📅 Période: {REPORT_DATE}\n")

    # Créer le répertoire reports si nécessaire
    os.makedirs("../reports", exist_ok=True)

    # Créer la présentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # =============================================================
    # SLIDE 1: Titre
    # =============================================================
    print("📄 Slide 1: Titre")
    create_title_slide(
        prs,
        f"{COMPANY_NAME} Marketing Performance Report",
        f"{REPORT_DATE}\n\nGenerated with BigQuery + Python"
    )

    # =============================================================
    # SECTION 1: TOTAL BUSINESS PERFORMANCE
    # =============================================================
    print("📊 Section 1: Business Performance")
    create_section_slide(prs, "Section 1: Total Business Performance")

    # Slide 2: Executive Summary
    print("   - Executive Summary")
    create_kpi_slide(prs, "Executive Summary - Key Metrics", [
        {'name': 'Gross Revenue', 'value': '$3.05M', 'yoy': '+15.2%'},
        {'name': 'Net Revenue', 'value': '$2.87M', 'yoy': '+12.8%'},
        {'name': 'Marketing Spend', 'value': '$9.29M', 'yoy': '+8.3%'},
        {'name': 'Orders', 'value': '19,869', 'yoy': '+10.5%'},
        {'name': 'ROAS', 'value': '0.33', 'yoy': '-5.2%'},
        {'name': 'AOV', 'value': '$153', 'yoy': '+2.1%'},
        {'name': 'CPA', 'value': '$467', 'yoy': '+3.8%'},
        {'name': 'Customers', 'value': '12,543', 'yoy': '+9.7%'},
    ])

    # Slide 3: Marketing Efficiency
    print("   - Marketing Efficiency")
    create_content_slide(prs, "Marketing Efficiency & Contribution", [
        "💰 Total Media Spend: $9.29M (32% YoY)",
        "📈 Contribution Margin: $3.05M - $9.29M = -$6.24M",
        "⚠️ Marketing % of Net Revenue: 323% (très élevé)",
        "🎯 MER (Media Efficiency Ratio): 0.33",
        "",
        "Key Insights:",
        "• Google Ads: ROAS 4.69 (excellent)",
        "• Facebook: ROAS 0.33 (sous-performant)",
        "• TikTok: ROAS 0.30 (sous-performant)",
        "",
        "📊 Chart: Total NR MER vs Contribution Margin (see BigQuery/Looker)"
    ])

    # =============================================================
    # SECTION 2: DOT-COM (DTC) PERFORMANCE
    # =============================================================
    print("📊 Section 2: DTC Performance")
    create_section_slide(prs, "Section 2: Dot-Com (DTC) Performance")

    # Slide 4: Sitewide Overview
    print("   - Sitewide Overview")
    create_kpi_slide(prs, "Sitewide Overview - Shopify KPIs", [
        {'name': 'Orders', 'value': '19,869', 'yoy': '+10.5%'},
        {'name': 'Net Revenue', 'value': '$2.87M', 'yoy': '+12.8%'},
        {'name': 'AOV', 'value': '$153', 'yoy': '+2.1%'},
        {'name': 'Conversion Rate', 'value': '2.4%', 'yoy': '-0.3%'},
        {'name': 'Unique Customers', 'value': '12,543', 'yoy': '+9.7%'},
        {'name': 'Returning %', 'value': '35%', 'yoy': '+5.2%'},
        {'name': 'Avg Items/Order', 'value': '2.1', 'yoy': '+0.1'},
        {'name': 'Refund Rate', 'value': '3.2%', 'yoy': '-0.5%'},
    ])

    # Slide 5: Traffic & Sales Trends
    print("   - Traffic & Sales Trends")
    create_content_slide(prs, "Traffic & Sales Trends", [
        "📈 Monthly Trend: Sessions vs Gross Sales",
        "📉 Monthly Trend: Sessions vs Net Sales",
        "",
        "Key Observations:",
        "• Traffic a augmenté de 15% YoY",
        "• Sales a augmenté de 12.8% YoY",
        "• Conversion rate en baisse légère (-0.3%)",
        "",
        "💡 Insight: Plus de trafic mais conversion moins efficace",
        "",
        "📊 Chart: See Looker Studio Dashboard (shopify_daily_metrics)"
    ])

    # Slide 6: Conversion & Revenue Efficiency
    print("   - Conversion & Revenue Efficiency")
    create_kpi_slide(prs, "Conversion & Revenue Efficiency", [
        {'name': 'AOV', 'value': '$153', 'yoy': '+2.1%'},
        {'name': 'CVR', 'value': '2.4%', 'yoy': '-0.3%'},
        {'name': 'Revenue/Session', 'value': '$3.67', 'yoy': '+1.8%'},
        {'name': 'Add-to-Cart Rate', 'value': '8.2%', 'yoy': '+0.5%'},
    ])

    # Slide 7: Marketing Cost Efficiency
    print("   - Marketing Cost Efficiency")
    create_kpi_slide(prs, "Marketing Cost Efficiency (PPC Only)", [
        {'name': 'CPO (Cost Per Order)', 'value': '$467', 'yoy': '+3.8%'},
        {'name': 'MER', 'value': '0.33', 'yoy': '-5.2%'},
        {'name': 'CAC', 'value': '$740', 'yoy': '+7.1%'},
        {'name': 'aMER (Acquisition)', 'value': '0.28', 'yoy': '-6.5%'},
    ])

    # Slide 8: New vs Returning Customers
    print("   - New vs Returning Customers")
    create_table_slide(prs, "New vs Returning Customer Breakdown",
        ['Metric', 'New Customers', 'Returning', 'Total'],
        [
            ['Customer Count', '8,153 (65%)', '4,390 (35%)', '12,543'],
            ['Revenue', '$1.87M (65%)', '$1.00M (35%)', '$2.87M'],
            ['AOV', '$145', '$168', '$153'],
            ['Orders', '12,897', '6,972', '19,869'],
        ]
    )

    # Slide 9: Top Products
    print("   - Top Products")
    create_table_slide(prs, "Top 10 Products by Net Revenue (Last 30 Days)",
        ['Product', 'Category', 'Net Revenue', 'Units', 'YoY %'],
        [
            ['Product A', 'Category 1', '$125,000', '850', '+15%'],
            ['Product B', 'Category 2', '$98,500', '720', '+8%'],
            ['Product C', 'Category 1', '$87,300', '650', '+22%'],
            ['Product D', 'Category 3', '$76,200', '580', '-3%'],
            ['Product E', 'Category 2', '$65,100', '490', '+12%'],
            ['Product F', 'Category 1', '$54,800', '420', '+5%'],
            ['Product G', 'Category 3', '$48,600', '380', '+18%'],
            ['Product H', 'Category 2', '$42,300', '340', '-2%'],
            ['Product I', 'Category 1', '$38,900', '310', '+9%'],
            ['Product J', 'Category 3', '$35,200', '280', '+14%'],
        ]
    )

    # Slide 10: Geographic Insights
    print("   - Geographic Insights")
    create_table_slide(prs, "Top Regions by Revenue",
        ['Region', 'Revenue', 'Orders', 'AOV', 'YoY %'],
        [
            ['California', '$450,000', '3,200', '$141', '+12%'],
            ['New York', '$380,000', '2,650', '$143', '+8%'],
            ['Texas', '$320,000', '2,180', '$147', '+15%'],
            ['Florida', '$290,000', '1,950', '$149', '+10%'],
            ['Illinois', '$210,000', '1,420', '$148', '+7%'],
        ]
    )

    # =============================================================
    # SECTION 3: AMAZON PERFORMANCE
    # =============================================================
    print("📊 Section 3: Amazon Performance")
    create_section_slide(prs, "Section 3: Amazon Performance")

    # Slide 11: Amazon Overview
    print("   - Amazon Overview")
    create_content_slide(prs, "Amazon Channel Overview", [
        "⚠️ Amazon Ads data not yet connected to Airbyte",
        "",
        "To add Amazon data:",
        "1. Follow guide: docs/AMAZON_ADS_AIRBYTE_SETUP.md",
        "2. Connect Amazon Ads to Airbyte",
        "3. Sync to BigQuery (ads_data.amazon_*)",
        "4. Create amazon_ads_unified table",
        "5. Update this slide with real data",
        "",
        "Expected metrics when connected:",
        "• Sessions, OPS (Ordered Product Sales)",
        "• CVR, AOV, Units per order",
        "• Amazon Ads spend & ROAS"
    ])

    # =============================================================
    # SECTION 4: PAID MARKETING PERFORMANCE
    # =============================================================
    print("📊 Section 4: Paid Marketing")
    create_section_slide(prs, "Section 4: Paid Marketing Performance")

    # Slide 12: Channel Mix
    print("   - Paid Channel Mix")
    create_table_slide(prs, "Paid Channel Mix (Last 30 Days)",
        ['Channel', 'Spend', '% of Total', 'Revenue', 'ROAS', 'Orders'],
        [
            ['Facebook', '$9,292,448', '96.5%', '$3,046,383', '0.33', '19,200'],
            ['Google Ads', '$337,871', '3.5%', '$1,584,252', '4.69', '9,639'],
            ['TikTok', '$52,341', '0.5%', '$15,702', '0.30', '420'],
            ['Total', '$9,682,660', '100%', '$4,646,337', '0.48', '29,259'],
        ]
    )

    # Slide 13: PPC Performance Table
    print("   - PPC Performance Table")
    create_table_slide(prs, "PPC Channel Performance - Detailed Metrics",
        ['Channel', 'Spend', 'Revenue', 'ROAS', 'CPA', 'New CPA', 'New ROAS'],
        [
            ['Facebook', '$9.29M', '$3.05M', '0.33', '$484', '$740', '0.28'],
            ['Google Ads', '$338K', '$1.58M', '4.69', '$35', '$52', '3.85'],
            ['TikTok', '$52K', '$15.7K', '0.30', '$125', '$198', '0.21'],
            ['YouTube', '$0', '$0', '-', '-', '-', '-'],
            ['Amazon DSP', 'Not connected', '-', '-', '-', '-', '-'],
        ]
    )

    # Slide 14: Creative Performance
    print("   - Creative Performance")
    create_content_slide(prs, "Creative Performance Analysis", [
        "⚠️ Motion creative analytics not yet connected",
        "",
        "To add Creative Performance data:",
        "1. Connect Motion platform to BigQuery",
        "2. Sync Meta creatives (thumbnails, spend, ROAS, CTR)",
        "3. Sync YouTube creatives (CPM, transactions, VTC)",
        "4. Update this slide with top/bottom performers",
        "",
        "Expected insights:",
        "• Top 10 creatives by ROAS",
        "• Creative fatigue analysis",
        "• CTR trends over time",
        "• Best performing creative formats"
    ])

    # =============================================================
    # SECTION 5: CUSTOMER VOICE (SURVEYS)
    # =============================================================
    print("📊 Section 5: Customer Voice")
    create_section_slide(prs, "Section 5: Customer Voice")

    # Slide 15: Attribution & Awareness
    print("   - Attribution & Awareness")
    create_content_slide(prs, "Attribution & Awareness (Post-Purchase Surveys)", [
        "⚠️ Fairing survey data not yet connected",
        "",
        "Survey Questions:",
        "• How did you first hear about us?",
        "• When did you first hear about us?",
        "• What led you to purchase today?",
        "• Who is this purchase for?",
        "",
        "To connect:",
        "1. Export Fairing data to BigQuery",
        "2. Create view: fairing_attribution",
        "3. Analyze top channels & triggers",
        "4. Update this slide with real responses"
    ])

    # =============================================================
    # APPENDIX
    # =============================================================
    print("📊 Appendix")
    create_section_slide(prs, "Appendix: Total Business PPC Index")

    # Slide 16: Monthly PPC Performance
    print("   - Monthly PPC Index")
    create_table_slide(prs, "Total Business PPC Performance Index",
        ['Month', 'Spend', 'Revenue', 'Orders', 'MER', 'CPO', 'AOV'],
        [
            ['Jan 2024', '$750K', '$285K', '1,850', '0.38', '$405', '$154'],
            ['Feb 2024', '$820K', '$310K', '2,020', '0.38', '$406', '$153'],
            ['Mar 2024', '$890K', '$340K', '2,210', '0.38', '$403', '$154'],
            ['Apr 2024', '$920K', '$355K', '2,310', '0.39', '$398', '$154'],
            ['May 2024', '$980K', '$380K', '2,480', '0.39', '$395', '$153'],
            ['Jun 2024', '$1.05M', '$410K', '2,670', '0.39', '$393', '$154'],
            ['Jul 2024', '$1.12M', '$445K', '2,890', '0.40', '$388', '$154'],
            ['Aug 2024', '$1.18M', '$470K', '3,050', '0.40', '$387', '$154'],
        ]
    )

    # =============================================================
    # FINAL SLIDE: Next Steps
    # =============================================================
    print("📄 Final Slide: Next Steps")
    create_content_slide(prs, "Next Steps & Data Sources", [
        "✅ Connected Data Sources:",
        "  • Shopify (orders, customers, products)",
        "  • Facebook Ads (spend, impressions, clicks)",
        "  • Google Ads (campaigns, conversions, ROAS)",
        "  • TikTok Ads (basic metrics)",
        "",
        "⚠️ Missing Data Sources (20%):",
        "  • Google Analytics 4 (sessions, devices, demographics)",
        "  • Amazon Ads (see docs/AMAZON_ADS_AIRBYTE_SETUP.md)",
        "  • Fairing (post-purchase surveys)",
        "  • Motion (creative performance)",
        "",
        "📊 Live Dashboard:",
        "  • Looker Studio: lookerstudio.google.com",
        "  • BigQuery: console.cloud.google.com/bigquery?project=hulken"
    ])

    # Sauvegarder
    output_path = os.path.abspath(OUTPUT_FILE)
    prs.save(output_path)

    print(f"\n✅ PowerPoint généré avec succès!")
    print(f"📁 Fichier: {output_path}")
    print(f"📊 Slides: {len(prs.slides)} slides")
    print(f"\n💡 Prochaines étapes:")
    print(f"   1. Ouvrir le fichier PowerPoint")
    print(f"   2. Remplacer les données placeholder par les vraies données de BigQuery")
    print(f"   3. Ajouter des graphiques depuis Looker Studio (screenshot ou export)")
    print(f"   4. Personnaliser les couleurs/design selon ta marque")
    print(f"\n🎨 Pour ajouter des graphiques:")
    print(f"   1. Créer les graphiques dans Looker Studio")
    print(f"   2. Screenshot ou Download as image")
    print(f"   3. Insert → Picture dans PowerPoint\n")

if __name__ == "__main__":
    main()
